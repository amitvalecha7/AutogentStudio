import os
import logging
from typing import Dict, List, Any, Optional
import mimetypes
from pathlib import Path

# File processing libraries
try:
    import PyPDF2
    from PyPDF2 import PdfReader
except ImportError:
    PyPDF2 = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from PIL import Image
except ImportError:
    Image = None

try:
    import eyed3
except ImportError:
    eyed3 = None

from services.ai_providers import get_embedding

def get_file_type(file_path: str) -> str:
    """Determine file type from path"""
    mime_type, _ = mimetypes.guess_type(file_path)
    extension = Path(file_path).suffix.lower()
    
    # Text documents
    if extension in ['.txt', '.md', '.rst']:
        return 'text'
    elif extension in ['.pdf']:
        return 'pdf'
    elif extension in ['.doc', '.docx', '.odt']:
        return 'document'
    elif extension in ['.xls', '.xlsx', '.csv', '.ods']:
        return 'spreadsheet'
    elif extension in ['.ppt', '.pptx', '.odp']:
        return 'presentation'
    
    # Images
    elif extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.svg']:
        return 'image'
    
    # Audio
    elif extension in ['.mp3', '.wav', '.ogg', '.m4a', '.flac']:
        return 'audio'
    
    # Video
    elif extension in ['.mp4', '.avi', '.mov', '.wmv', '.webm']:
        return 'video'
    
    # Code files
    elif extension in ['.py', '.js', '.html', '.css', '.java', '.cpp', '.c', '.h']:
        return 'code'
    
    # Data files
    elif extension in ['.json', '.xml', '.yaml', '.yml']:
        return 'data'
    
    else:
        return 'unknown'

def extract_text(file_path: str, file_type: str) -> str:
    """Extract text content from various file types"""
    try:
        if file_type == 'text' or file_type == 'code' or file_type == 'data':
            return extract_text_file(file_path)
        elif file_type == 'pdf':
            return extract_pdf_text(file_path)
        elif file_type == 'document':
            return extract_document_text(file_path)
        elif file_type == 'spreadsheet':
            return extract_spreadsheet_text(file_path)
        elif file_type == 'presentation':
            return extract_presentation_text(file_path)
        elif file_type == 'image':
            return extract_image_metadata(file_path)
        elif file_type == 'audio':
            return extract_audio_metadata(file_path)
        elif file_type == 'video':
            return extract_video_metadata(file_path)
        else:
            return f"File type '{file_type}' not supported for text extraction"
            
    except Exception as e:
        logging.error(f"Error extracting text from {file_path}: {e}")
        return f"Error extracting text: {str(e)}"

def extract_text_file(file_path: str) -> str:
    """Extract text from plain text files"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Try with different encoding
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            return f"Error reading file with encoding: {str(e)}"
    except Exception as e:
        return f"Error reading text file: {str(e)}"

def extract_pdf_text(file_path: str) -> str:
    """Extract text from PDF files"""
    if not PyPDF2:
        return "PyPDF2 not installed - PDF text extraction not available"
    
    try:
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                try:
                    text += page.extract_text() + "\n"
                except Exception as e:
                    text += f"[Error extracting page: {str(e)}]\n"
        return text.strip()
    except Exception as e:
        return f"Error extracting PDF text: {str(e)}"

def extract_document_text(file_path: str) -> str:
    """Extract text from Word documents"""
    extension = Path(file_path).suffix.lower()
    
    if extension == '.docx':
        if not DocxDocument:
            return "python-docx not installed - DOCX text extraction not available"
        
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except Exception as e:
            return f"Error extracting DOCX text: {str(e)}"
    
    else:
        return f"Document format '{extension}' not supported"

def extract_spreadsheet_text(file_path: str) -> str:
    """Extract text from spreadsheet files"""
    if not pd:
        return "pandas not installed - Spreadsheet text extraction not available"
    
    extension = Path(file_path).suffix.lower()
    
    try:
        if extension == '.csv':
            df = pd.read_csv(file_path)
        elif extension in ['.xls', '.xlsx']:
            df = pd.read_excel(file_path)
        else:
            return f"Spreadsheet format '{extension}' not supported"
        
        # Convert to text representation
        text = f"Spreadsheet content ({df.shape[0]} rows, {df.shape[1]} columns):\n"
        text += f"Columns: {', '.join(df.columns.astype(str))}\n\n"
        
        # Add first few rows as sample
        sample_rows = min(5, len(df))
        text += f"Sample data (first {sample_rows} rows):\n"
        text += df.head(sample_rows).to_string(index=False)
        
        return text
        
    except Exception as e:
        return f"Error extracting spreadsheet text: {str(e)}"

def extract_presentation_text(file_path: str) -> str:
    """Extract text from presentation files"""
    if not Presentation:
        return "python-pptx not installed - Presentation text extraction not available"
    
    extension = Path(file_path).suffix.lower()
    
    if extension == '.pptx':
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            
            return text.strip()
            
        except Exception as e:
            return f"Error extracting presentation text: {str(e)}"
    
    else:
        return f"Presentation format '{extension}' not supported"

def extract_image_metadata(file_path: str) -> str:
    """Extract metadata from image files"""
    if not Image:
        return "Pillow not installed - Image metadata extraction not available"
    
    try:
        with Image.open(file_path) as img:
            text = f"Image metadata:\n"
            text += f"Format: {img.format}\n"
            text += f"Mode: {img.mode}\n"
            text += f"Size: {img.size[0]} x {img.size[1]} pixels\n"
            
            # Extract EXIF data if available
            if hasattr(img, '_getexif') and img._getexif():
                exif = img._getexif()
                if exif:
                    text += "EXIF data available\n"
            
            return text
            
    except Exception as e:
        return f"Error extracting image metadata: {str(e)}"

def extract_audio_metadata(file_path: str) -> str:
    """Extract metadata from audio files"""
    if not eyed3:
        return "eyed3 not installed - Audio metadata extraction not available"
    
    try:
        audiofile = eyed3.load(file_path)
        if audiofile.tag:
            text = f"Audio metadata:\n"
            text += f"Title: {audiofile.tag.title or 'Unknown'}\n"
            text += f"Artist: {audiofile.tag.artist or 'Unknown'}\n"
            text += f"Album: {audiofile.tag.album or 'Unknown'}\n"
            text += f"Duration: {audiofile.info.time_secs if audiofile.info else 'Unknown'} seconds\n"
            return text
        else:
            return "No audio metadata found"
            
    except Exception as e:
        return f"Error extracting audio metadata: {str(e)}"

def extract_video_metadata(file_path: str) -> str:
    """Extract metadata from video files"""
    # Basic file info since video processing libraries are complex
    try:
        file_size = os.path.getsize(file_path)
        text = f"Video file metadata:\n"
        text += f"File size: {file_size} bytes\n"
        text += f"File path: {file_path}\n"
        return text
        
    except Exception as e:
        return f"Error extracting video metadata: {str(e)}"

def process_file(file_id: str) -> Dict[str, Any]:
    """Process file and return processing results"""
    from models import File, FileChunk, KnowledgeBase
    from app import db
    
    try:
        # Get file record
        file_record = File.query.get(file_id)
        if not file_record:
            return {'success': False, 'error': 'File not found'}
        
        # Update status
        file_record.processing_status = 'processing'
        db.session.commit()
        
        # Determine file type
        file_type = get_file_type(file_record.file_path)
        
        # Extract text content
        text_content = extract_text(file_record.file_path, file_type)
        
        result = {
            'success': True,
            'file_type': file_type,
            'text_content': text_content,
            'chunks_created': 0
        }
        
        # If file belongs to a knowledge base, create chunks
        if file_record.knowledge_base_id and text_content:
            kb = KnowledgeBase.query.get(file_record.knowledge_base_id)
            if kb:
                chunks = chunk_text(text_content, kb.chunk_size, kb.chunk_overlap)
                
                # Create file chunks with embeddings
                for i, chunk_content in enumerate(chunks):
                    try:
                        # Generate embedding
                        embedding = get_embedding(chunk_content, kb.embedding_model)
                        
                        # Convert embedding to bytes for storage
                        import numpy as np
                        embedding_bytes = np.array(embedding, dtype=np.float32).tobytes()
                        
                        chunk = FileChunk(
                            file_id=file_record.id,
                            knowledge_base_id=kb.id,
                            chunk_index=i,
                            content=chunk_content,
                            embedding=embedding_bytes,
                            embedding_model=kb.embedding_model,
                            metadata={
                                'file_name': file_record.original_filename,
                                'file_type': file_type,
                                'chunk_size': len(chunk_content)
                            }
                        )
                        db.session.add(chunk)
                        
                    except Exception as e:
                        logging.error(f"Error creating chunk {i}: {e}")
                        continue
                
                result['chunks_created'] = len(chunks)
        
        # Update file record
        file_record.processing_status = 'completed'
        file_record.is_processed = True
        file_record.metadata = {
            'file_type': file_type,
            'text_length': len(text_content),
            'chunks_created': result['chunks_created']
        }
        
        db.session.commit()
        
        return result
        
    except Exception as e:
        logging.error(f"Error processing file {file_id}: {e}")
        
        # Update file status to failed
        try:
            file_record = File.query.get(file_id)
            if file_record:
                file_record.processing_status = 'failed'
                db.session.commit()
        except:
            pass
        
        return {'success': False, 'error': str(e)}

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """Split text into overlapping chunks"""
    if not text:
        return []
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + chunk_size
        
        # If this isn't the last chunk, try to break at a sentence or word boundary
        if end < len(text):
            # Look for sentence boundary (period, exclamation, question mark)
            sentence_end = max(
                text.rfind('.', start, end),
                text.rfind('!', start, end),
                text.rfind('?', start, end)
            )
            
            if sentence_end > start + chunk_size // 2:
                end = sentence_end + 1
            else:
                # Look for word boundary (space)
                word_end = text.rfind(' ', start, end)
                if word_end > start + chunk_size // 2:
                    end = word_end
        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        
        # Move start position with overlap
        start = end - overlap
        
        # Ensure we don't get stuck in infinite loop
        if start >= len(text):
            break
    
    return chunks

def generate_embeddings(text: str, model: str = 'text-embedding-3-small') -> List[float]:
    """Generate embeddings for text"""
    try:
        return get_embedding(text, model)
    except Exception as e:
        logging.error(f"Error generating embeddings: {e}")
        raise

def validate_file_upload(file, max_size: int = 16 * 1024 * 1024) -> Dict[str, Any]:
    """Validate uploaded file"""
    if not file or file.filename == '':
        return {'valid': False, 'error': 'No file provided'}
    
    # Check file size
    file.seek(0, 2)  # Seek to end
    file_size = file.tell()
    file.seek(0)  # Reset to beginning
    
    if file_size > max_size:
        return {'valid': False, 'error': f'File too large. Maximum size: {max_size // (1024*1024)}MB'}
    
    # Check file extension
    allowed_extensions = {
        'txt', 'pdf', 'doc', 'docx', 'rtf', 'odt',
        'ppt', 'pptx', 'odp',
        'xls', 'xlsx', 'ods', 'csv',
        'jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp', 'svg',
        'mp3', 'wav', 'ogg', 'm4a', 'flac',
        'mp4', 'avi', 'mov', 'wmv', 'webm',
        'html', 'htm', 'xml', 'json', 'yaml', 'yml',
        'py', 'js', 'css', 'java', 'cpp', 'c', 'h',
        'md', 'rst', 'tex'
    }
    
    if '.' in file.filename:
        extension = file.filename.rsplit('.', 1)[1].lower()
        if extension not in allowed_extensions:
            return {'valid': False, 'error': f'File type .{extension} not supported'}
    else:
        return {'valid': False, 'error': 'File must have an extension'}
    
    return {'valid': True, 'file_size': file_size}

def get_file_preview(file_path: str, max_length: int = 500) -> str:
    """Get preview of file content"""
    try:
        file_type = get_file_type(file_path)
        full_text = extract_text(file_path, file_type)
        
        if len(full_text) <= max_length:
            return full_text
        else:
            return full_text[:max_length] + "..."
            
    except Exception as e:
        logging.error(f"Error generating file preview: {e}")
        return f"Error generating preview: {str(e)}"

def cleanup_temp_files(temp_dir: str = 'temp') -> None:
    """Clean up temporary files older than 1 hour"""
    import time
    
    try:
        if not os.path.exists(temp_dir):
            return
        
        current_time = time.time()
        one_hour = 3600  # seconds
        
        for filename in os.listdir(temp_dir):
            file_path = os.path.join(temp_dir, filename)
            if os.path.isfile(file_path):
                file_age = current_time - os.path.getmtime(file_path)
                if file_age > one_hour:
                    try:
                        os.remove(file_path)
                        logging.info(f"Cleaned up temporary file: {file_path}")
                    except Exception as e:
                        logging.error(f"Error removing temporary file {file_path}: {e}")
                        
    except Exception as e:
        logging.error(f"Error during temp file cleanup: {e}")
