import os
import logging
import mimetypes
from pathlib import Path
from app import db
from models import File, FileChunk
from services.vector_service import VectorService

class FileService:
    def __init__(self):
        self.vector_service = VectorService()
        self.supported_text_formats = {
            'text/plain': self._extract_text_plain,
            'application/pdf': self._extract_pdf,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._extract_docx,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._extract_xlsx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._extract_pptx,
            'text/markdown': self._extract_markdown,
            'text/csv': self._extract_csv,
            'application/json': self._extract_json
        }
    
    def process_file_for_rag(self, file_id):
        """Process file for RAG - extract text and create embeddings"""
        try:
            file_record = File.query.get(file_id)
            if not file_record:
                raise Exception(f"File {file_id} not found")
            
            # Extract text content
            text_content = self._extract_text_from_file(file_record)
            if not text_content:
                file_record.processing_status = 'failed'
                file_record.is_processed = False
                db.session.commit()
                return
            
            # Create chunks
            chunks = self._create_chunks(text_content)
            
            # Save chunks and create embeddings
            for i, chunk in enumerate(chunks):
                file_chunk = FileChunk(
                    file_id=file_id,
                    chunk_index=i,
                    content=chunk
                )
                
                # Create embedding
                try:
                    embedding = self.vector_service.create_embedding(chunk)
                    file_chunk.embedding = embedding
                except Exception as e:
                    logging.error(f"Error creating embedding for chunk {i}: {str(e)}")
                
                db.session.add(file_chunk)
            
            file_record.processing_status = 'completed'
            file_record.is_processed = True
            db.session.commit()
            
            logging.info(f"File {file_id} processed successfully with {len(chunks)} chunks")
            
        except Exception as e:
            logging.error(f"Error processing file {file_id}: {str(e)}")
            file_record = File.query.get(file_id)
            if file_record:
                file_record.processing_status = 'failed'
                file_record.is_processed = False
                db.session.commit()
    
    def _extract_text_from_file(self, file_record):
        """Extract text content from file based on MIME type"""
        mime_type = file_record.mime_type or mimetypes.guess_type(file_record.file_path)[0]
        
        if mime_type in self.supported_text_formats:
            return self.supported_text_formats[mime_type](file_record.file_path)
        else:
            logging.warning(f"Unsupported file type: {mime_type}")
            return None
    
    def _extract_text_plain(self, file_path):
        """Extract text from plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logging.error(f"Error reading text file: {str(e)}")
            return None
    
    def _extract_pdf(self, file_path):
        """Extract text from PDF file"""
        try:
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages:
                    text += page.extract_text()
                return text
        except ImportError:
            logging.error("PyPDF2 not installed")
            return None
        except Exception as e:
            logging.error(f"Error reading PDF file: {str(e)}")
            return None
    
    def _extract_docx(self, file_path):
        """Extract text from DOCX file"""
        try:
            from docx import Document
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except ImportError:
            logging.error("python-docx not installed")
            return None
        except Exception as e:
            logging.error(f"Error reading DOCX file: {str(e)}")
            return None
    
    def _extract_xlsx(self, file_path):
        """Extract text from XLSX file"""
        try:
            import pandas as pd
            df = pd.read_excel(file_path)
            return df.to_string()
        except ImportError:
            logging.error("pandas not installed")
            return None
        except Exception as e:
            logging.error(f"Error reading XLSX file: {str(e)}")
            return None
    
    def _extract_pptx(self, file_path):
        """Extract text from PPTX file"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except ImportError:
            logging.error("python-pptx not installed")
            return None
        except Exception as e:
            logging.error(f"Error reading PPTX file: {str(e)}")
            return None
    
    def _extract_markdown(self, file_path):
        """Extract text from Markdown file"""
        return self._extract_text_plain(file_path)
    
    def _extract_csv(self, file_path):
        """Extract text from CSV file"""
        try:
            import pandas as pd
            df = pd.read_csv(file_path)
            return df.to_string()
        except ImportError:
            logging.error("pandas not installed")
            return None
        except Exception as e:
            logging.error(f"Error reading CSV file: {str(e)}")
            return None
    
    def _extract_json(self, file_path):
        """Extract text from JSON file"""
        try:
            import json
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return json.dumps(data, indent=2)
        except Exception as e:
            logging.error(f"Error reading JSON file: {str(e)}")
            return None
    
    def _create_chunks(self, text, chunk_size=1000, overlap=200):
        """Create text chunks with overlap"""
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            chunks.append(chunk)
            start = end - overlap
        
        return chunks
    
    def search_files(self, query, user_id, limit=10):
        """Search files using vector similarity"""
        try:
            # Create embedding for query
            query_embedding = self.vector_service.create_embedding(query)
            
            # Search for similar chunks
            similar_chunks = self.vector_service.search_similar_chunks(
                query_embedding, user_id, limit
            )
            
            return similar_chunks
        except Exception as e:
            logging.error(f"Error searching files: {str(e)}")
            return []
